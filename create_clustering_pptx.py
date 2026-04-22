"""
Generate NON-OFS ILEC Master Clustering Data Dictionary PPTX deck.
Run with system Python 3.11 (has python-pptx).
Output: NON_OFS_Clustering_DataDictionary.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──
VZ_RED = RGBColor(0xEE, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = r"C:\Users\v267429\Downloads\AI_Sessions"
OUT_FILE = os.path.join(OUT_DIR, "NON_OFS_Clustering_DataDictionary.pptx")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_red_bar(slide, top=Inches(0), height=Inches(1.15)):
    """Full-width Verizon red banner at top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = VZ_RED
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
                font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    tf.vertical_anchor = anchor
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=14, color=DARK_GRAY, font_name="Calibri",
                          line_spacing=1.15, bold_first=False, bullet=False):
    """lines is a list of strings. Returns the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet:
            p.text = f"\u2022  {line}"
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=VZ_RED, bg_color=LIGHT_GRAY):
    """Big number card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Number
    add_textbox(slide, left, top + Inches(0.15), width, Inches(0.6),
                text=number, font_size=32, bold=True, color=num_color,
                align=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left, top + Inches(0.75), width, Inches(0.5),
                text=label, font_size=13, bold=False, color=MED_GRAY,
                align=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, height, headers, rows,
              header_bg=VZ_RED, header_fg=WHITE, col_widths=None):
    """Add a table with styled header row."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = header_fg
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl


def add_quadrant_cell(slide, left, top, width, height, label, detail,
                      bg_color, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(2)
    shape.shadow.inherit = False
    # Label
    add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.45),
                text=label, font_size=16, bold=True, color=text_color,
                align=PP_ALIGN.LEFT)
    # Detail
    add_multiline_textbox(slide, left + Inches(0.15), top + Inches(0.55),
                          width - Inches(0.3), height - Inches(0.7),
                          detail, font_size=11, color=text_color)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s1 = add_blank_slide()
add_red_bar(s1, top=Inches(0), height=Inches(7.5))  # full red background

# Decorative white stripe
stripe = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), SLIDE_W, Inches(2.5))
stripe.fill.solid()
stripe.fill.fore_color.rgb = WHITE
stripe.line.fill.background()

add_textbox(s1, Inches(0.8), Inches(2.65), Inches(11), Inches(0.9),
            "NON-OFS ILEC Master Clustering", font_size=40, bold=True,
            color=VZ_RED, align=PP_ALIGN.LEFT)

add_textbox(s1, Inches(0.8), Inches(3.55), Inches(11), Inches(0.6),
            "Data Dictionary & Build Prioritization Methodology", font_size=22,
            bold=False, color=MED_GRAY, align=PP_ALIGN.LEFT)

add_textbox(s1, Inches(0.8), Inches(4.2), Inches(4), Inches(0.4),
            "April 22, 2026", font_size=16, bold=False, color=MED_GRAY,
            align=PP_ALIGN.LEFT)

# Verizon branding bottom
add_textbox(s1, Inches(0.8), Inches(6.5), Inches(4), Inches(0.5),
            "Verizon  |  Broadband Strategy & Engineering", font_size=14,
            bold=False, color=WHITE, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ════════════════════════════════════════════════════════════════════
s2 = add_blank_slide()
add_red_bar(s2)
add_textbox(s2, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Executive Summary", font_size=30, bold=True, color=WHITE)

# Three pillars
pillars = [
    ("What", "12.26M NON-OFS ILEC addresses clustered into 68,071 build-ready hubs across 2,078 wire centers, with financial quality scoring and obligation tagging."),
    ("How", "Road-graph geospatial clustering using TIGER road networks. 6-phase algorithm: score, snap-to-road, build adjacency graph, seed & grow hubs, enforce size constraints, finalize. Financial quality drives hub formation; obligations are tagged post-formation for stability."),
    ("Output", "Oracle tables (CLUSTER_ASSIGNMENT_MASTER at 12.26M rows, CLUSTER_SUMMARY_MASTER at 68K rows) + ArcGIS File GDB (NONOFS_Master.gdb) + 54 MB GeoJSON for Tableau, ArcGIS, and dashboards. Eisenhower 2x2 matrix prioritizes every cluster for build sequencing."),
]

for i, (title, body) in enumerate(pillars):
    col_left = Inches(0.5) + Inches(i * 4.1)
    # Title box
    box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_left, Inches(1.5), Inches(3.9), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = VZ_RED
    box.line.fill.background()
    add_textbox(s2, col_left + Inches(0.15), Inches(1.52), Inches(3.6), Inches(0.5),
                text=title, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Body
    add_textbox(s2, col_left + Inches(0.15), Inches(2.2), Inches(3.7), Inches(4.5),
                text=body, font_size=13, color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Statistics
# ════════════════════════════════════════════════════════════════════
s3 = add_blank_slide()
add_red_bar(s3)
add_textbox(s3, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Key Statistics", font_size=30, bold=True, color=WHITE)

stats = [
    ("12.26M", "Addresses"),
    ("15.9M", "Total Units"),
    ("2,078", "Wire Centers"),
    ("68,071", "Clusters"),
    ("26.6%", "Obligated"),
]

card_w = Inches(2.2)
card_h = Inches(1.35)
gap = Inches(0.25)
total_w = len(stats) * card_w.inches + (len(stats) - 1) * gap.inches
start_x = (13.333 - total_w) / 2

for i, (num, lbl) in enumerate(stats):
    x = Inches(start_x + i * (card_w.inches + gap.inches))
    add_stat_box(s3, x, Inches(2.0), card_w, card_h, num, lbl)

# Additional detail rows
detail_lines = [
    "Avg addresses per cluster: 180  |  Avg units per cluster: 233",
    "Clusters with obligations: 25,333 (37.2%)  |  GeoJSON polygon features: 65,386",
    "ArcGIS File GDB: NONOFS_Master.gdb (~2 GB)  |  Oracle: TABLEAU_USER schema",
]
add_multiline_textbox(s3, Inches(1.0), Inches(4.0), Inches(11), Inches(2.5),
                      detail_lines, font_size=14, color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Formation Scoring
# ════════════════════════════════════════════════════════════════════
s4 = add_blank_slide()
add_red_bar(s4)
add_textbox(s4, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Formation Scoring — How Clusters Are Built", font_size=30,
            bold=True, color=WHITE)

# Subtitle
add_textbox(s4, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
            "Each address receives a 0-100 financial quality score that drives hub seed priority and growth order.",
            font_size=14, color=MED_GRAY)

# Table
headers4 = ["Component", "Weight", "What It Measures", "Normalization"]
rows4 = [
    ["IRR V2", "40%", "15-year internal rate of return from V2 ramp model", "Min-max within wire center, higher = better"],
    ["Copper Salvage", "25%", "Estimated recoverable value of existing copper plant", "copper_circuits x $200/circuit, normalized to WC max"],
    ["BI Priority Rank", "20%", "BI composite rank (dispatch, sales, geo synergy)", "Inverted min-max (lower rank = higher score)"],
    ["Unit Density", "15%", "MDU/MTU density bonus", "min(units / 10, 1.0) - caps at 10+ units"],
]
add_table(s4, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.2),
          headers4, rows4,
          col_widths=[Inches(1.8), Inches(0.9), Inches(4.5), Inches(5.1)])

# Why obligations excluded
box4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.6))
box4.fill.solid()
box4.fill.fore_color.rgb = LIGHT_GRAY
box4.line.fill.background()

add_textbox(s4, Inches(0.7), Inches(4.7), Inches(11.9), Inches(0.4),
            "Why are obligations excluded from formation?", font_size=14,
            bold=True, color=VZ_RED)

add_textbox(s4, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.9),
            "Obligations change frequently (new mandates, date shifts, scope changes) but cluster "
            "boundaries should be stable. Separating obligation tagging from formation means monthly "
            "score refreshes don't require re-clustering. The 6-phase algorithm uses only financial "
            "quality to build hubs; obligations are applied post-formation as tags.",
            font_size=12, color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Obligation Buckets
# ════════════════════════════════════════════════════════════════════
s5 = add_blank_slide()
add_red_bar(s5)
add_textbox(s5, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Obligation Buckets", font_size=30, bold=True, color=WHITE)

add_textbox(s5, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
            "Each address is tagged with exactly one obligation bucket. Priority-ordered \u2014 first match wins.",
            font_size=14, color=MED_GRAY)

headers5 = ["Priority", "Bucket", "Definition", "Count", "%"]
rows5 = [
    ["1", "COP_2026_OBLIG", "Planned copper recycling area, start 2026. Fiber must precede copper retirement.", "496,908", "4.1%"],
    ["2", "COP_2027_OBLIG", "Planned copper recycling area, start 2027.", "386,214", "3.1%"],
    ["3", "COP_FUTURE_OBLIG", "Copper recycling planned, start 2028+.", "14,385", "0.1%"],
    ["4", "SBB_OBLIG", "State broadband obligation area (grant or regulatory commitment).", "2,344,904", "19.1%"],
    ["5", "NSI_OBLIG", "National Security Interest inquiry match via NTAS bridge.", "20,280", "0.2%"],
    ["6", "LFA_OBLIG", "Local Franchise Authority obligation area.", "0", "0.0%"],
    ["7", "DISCRETIONARY", "No external obligation \u2014 purely economic decision.", "8,999,471", "73.4%"],
]
add_table(s5, Inches(0.5), Inches(1.95), Inches(12.3), Inches(3.6),
          headers5, rows5,
          col_widths=[Inches(0.8), Inches(1.9), Inches(6.0), Inches(1.5), Inches(0.8)])

# Key insight
box5 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9))
box5.fill.solid()
box5.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
box5.line.color.rgb = VZ_RED
box5.line.width = Pt(1.5)

add_textbox(s5, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7),
            "Key Insight:  3.26M addresses (26.6%) carry obligations \u2014 these represent mandatory builds "
            "driven by copper retirement, state broadband commitments, or national security inquiries. "
            "The remaining 73.4% are discretionary and prioritized by financial return.",
            font_size=12, bold=False, color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Eisenhower Matrix
# ════════════════════════════════════════════════════════════════════
s6 = add_blank_slide()
add_red_bar(s6)
add_textbox(s6, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Eisenhower Matrix \u2014 Build Prioritization", font_size=30,
            bold=True, color=WHITE)

# Axis labels
add_textbox(s6, Inches(1.1), Inches(1.4), Inches(1.2), Inches(0.5),
            "URGENCY \u2192", font_size=13, bold=True, color=MED_GRAY,
            align=PP_ALIGN.CENTER)

# Column headers
add_textbox(s6, Inches(3.0), Inches(1.4), Inches(4.5), Inches(0.4),
            "Urgency >= 50", font_size=13, bold=True, color=VZ_RED,
            align=PP_ALIGN.CENTER)
add_textbox(s6, Inches(7.8), Inches(1.4), Inches(4.5), Inches(0.4),
            "Urgency < 50", font_size=13, bold=True, color=MED_GRAY,
            align=PP_ALIGN.CENTER)

# Row labels
add_textbox(s6, Inches(0.3), Inches(2.3), Inches(2.5), Inches(0.4),
            "Value >= 50", font_size=13, bold=True, color=VZ_RED,
            align=PP_ALIGN.RIGHT)
add_textbox(s6, Inches(0.3), Inches(4.7), Inches(2.5), Inches(0.4),
            "Value < 50", font_size=13, bold=True, color=MED_GRAY,
            align=PP_ALIGN.RIGHT)

quad_w = Inches(4.5)
quad_h = Inches(2.15)

# Q1 — Do First (top-left)
add_quadrant_cell(s6, Inches(3.0), Inches(1.9), quad_w, quad_h,
                  "Q1 \u2014 Do First",
                  ["6,651 clusters  |  1.1M units",
                   "",
                   "Build immediately: high financial",
                   "return AND obligation / operational",
                   "pressure. Copper retirement +",
                   "strong IRR = top priority."],
                  RGBColor(0xCC, 0x00, 0x00))

# Q2 — Schedule (top-right)
add_quadrant_cell(s6, Inches(7.8), Inches(1.9), quad_w, quad_h,
                  "Q2 \u2014 Schedule",
                  ["44,525 clusters  |  11.1M units",
                   "",
                   "Economically attractive but no",
                   "urgency. Schedule into build plan",
                   "in ROI-descending order. Largest",
                   "bucket by far."],
                  RGBColor(0x00, 0x66, 0xCC))

# Q3 — Must Do (bottom-left)
add_quadrant_cell(s6, Inches(3.0), Inches(4.2), quad_w, quad_h,
                  "Q3 \u2014 Must Do",
                  ["1,039 clusters  |  98K units",
                   "",
                   "Obligation-driven but lower returns.",
                   "Build to meet regulatory commitments",
                   "even though IRR is below threshold."],
                  RGBColor(0xFF, 0x99, 0x00))

# Q4 — Deprioritize (bottom-right)
add_quadrant_cell(s6, Inches(7.8), Inches(4.2), quad_w, quad_h,
                  "Q4 \u2014 Deprioritize",
                  ["15,856 clusters  |  3.5M units",
                   "",
                   "Low return, no obligation. Defer",
                   "or evaluate selectively. May become",
                   "viable with cost reductions or new",
                   "obligations."],
                  RGBColor(0x88, 0x88, 0x88))

# VALUE axis label (vertical, left side)
add_textbox(s6, Inches(0.3), Inches(3.4), Inches(2.5), Inches(0.5),
            "VALUE \u2191", font_size=13, bold=True, color=MED_GRAY,
            align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Urgency & Value Scoring
# ════════════════════════════════════════════════════════════════════
s7 = add_blank_slide()
add_red_bar(s7)
add_textbox(s7, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Urgency & Value Scoring (0\u2013100)", font_size=30,
            bold=True, color=WHITE)

# Urgency section
add_textbox(s7, Inches(0.5), Inches(1.4), Inches(6), Inches(0.45),
            "Urgency Score \u2014 \"How pressing is the need to build?\"",
            font_size=16, bold=True, color=VZ_RED)

urg_headers = ["Component", "Weight", "What It Measures", "Scoring"]
urg_rows = [
    ["Obligation Tier", "50%", "Highest obligation bucket in cluster",
     "COP_2026=100, COP_2027=80, COP_FUTURE=60, SBB=70, NSI=50, LFA=40, DISC=0"],
    ["Copper Imminence", "30%", "% of addresses with copper retirement 2026-2027",
     "Linear: 0% \u2192 0, 100% \u2192 100"],
    ["Dispatch Frequency", "20%", "Avg truck rolls per address (past 12 months)",
     "min(avg_dispatch x 20, 100) \u2014 5+/yr = 100"],
]
add_table(s7, Inches(0.5), Inches(1.95), Inches(12.3), Inches(1.7),
          urg_headers, urg_rows,
          col_widths=[Inches(1.8), Inches(0.9), Inches(3.8), Inches(5.8)])

# Value section
add_textbox(s7, Inches(0.5), Inches(4.0), Inches(6), Inches(0.45),
            "Value Score \u2014 \"How financially attractive is this cluster?\"",
            font_size=16, bold=True, color=VZ_RED)

val_headers = ["Component", "Weight", "What It Measures", "Scoring"]
val_rows = [
    ["IRR V2", "50%", "Average 15-year IRR across cluster",
     "0% \u2192 0, 15% \u2192 50, 30%+ \u2192 100 (linear)"],
    ["Terminal Penetration", "25%", "Average projected FiOS take rate",
     "min(pen x 200, 100) \u2014 50%+ terminal = 100"],
    ["Revenue Potential", "25%", "Average annual EBITDA per address",
     "min(ebitda / $500, 100) \u2014 $500+/yr = 100"],
]
add_table(s7, Inches(0.5), Inches(4.55), Inches(12.3), Inches(1.7),
          val_headers, val_rows,
          col_widths=[Inches(1.8), Inches(0.9), Inches(3.8), Inches(5.8)])

# Note box
note7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.65))
note7.fill.solid()
note7.fill.fore_color.rgb = LIGHT_GRAY
note7.line.fill.background()
add_textbox(s7, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
            "Clusters scoring >= 50 on both axes are assigned Q1 (Do First). "
            "Urgency and Value are independent axes \u2014 a cluster can be high-value but low-urgency (Q2) or vice versa (Q3).",
            font_size=11, color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Output Tables
# ════════════════════════════════════════════════════════════════════
s8 = add_blank_slide()
add_red_bar(s8)
add_textbox(s8, Inches(0.6), Inches(0.2), Inches(10), Inches(0.8),
            "Output Tables \u2014 Oracle & File Outputs", font_size=30,
            bold=True, color=WHITE)

# CLUSTER_ASSIGNMENT_MASTER section
add_textbox(s8, Inches(0.5), Inches(1.4), Inches(12), Inches(0.45),
            "CLUSTER_ASSIGNMENT_MASTER  \u2014  12,262,162 rows (address-level)",
            font_size=15, bold=True, color=VZ_RED)

add_textbox(s8, Inches(0.5), Inches(1.85), Inches(12), Inches(0.4),
            "Key: (LOCUS_ADDRESS_ID, RUN_DATE)  |  Location: TABLEAU_USER schema",
            font_size=11, color=MED_GRAY)

cam_headers = ["Column", "Description"]
cam_rows = [
    ["LOCUS_ADDRESS_ID", "Unique address identifier"],
    ["CLUSTER_ID", "Hub assignment (format: {CLLI}_H{nnn})"],
    ["CLLI / REGION / STATE", "Wire center, region, state"],
    ["AUI / NO_OF_UNITS", "Address use indicator + NTAS unit count"],
    ["CPO_NTAS / CPO_PRED / TOTAL_CAPEX", "Cost metrics: per-unit, per-address, total"],
    ["COMPUTED_IRR / PEN_TERMINAL", "15-yr IRR (V2 ramp) and terminal penetration rate"],
    ["OBLIGATION_BUCKET", "COP_2026, COP_2027, SBB, NSI, DISCRETIONARY, etc."],
    ["FORMATION_SCORE", "Financial quality score (0-100) used during clustering"],
]
add_table(s8, Inches(0.5), Inches(2.25), Inches(12.3), Inches(2.15),
          cam_headers, cam_rows,
          col_widths=[Inches(3.5), Inches(8.8)])

# CLUSTER_SUMMARY_MASTER section
add_textbox(s8, Inches(0.5), Inches(4.6), Inches(12), Inches(0.45),
            "CLUSTER_SUMMARY_MASTER  \u2014  68,071 rows (cluster-level)",
            font_size=15, bold=True, color=VZ_RED)

add_textbox(s8, Inches(0.5), Inches(5.05), Inches(12), Inches(0.4),
            "Key: (CLUSTER_ID, RUN_DATE)  |  Location: TABLEAU_USER schema",
            font_size=11, color=MED_GRAY)

csm_headers = ["Column", "Description"]
csm_rows = [
    ["CLUSTER_ID / CLLI / REGION", "Cluster identity and geography"],
    ["TOTAL_UNITS / TOTAL_ADDRS / TOTAL_CAPEX", "Aggregate size and cost"],
    ["AVG_CPP / MEDIAN_IRR", "Cost Per Premises and median return"],
    ["URGENCY_SCORE / VALUE_SCORE", "Eisenhower axes (0-100 each)"],
    ["BUILD_PRIORITY_TIER", "Q1_Do_First, Q2_Schedule, Q3_Must_Do, Q4_Deprioritize"],
    ["OBLIGATION_FILL", "Pipe-delimited bucket breakdown (e.g., COP_2026:42|SBB:118)"],
]
add_table(s8, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7),
          csm_headers, csm_rows,
          col_widths=[Inches(3.5), Inches(8.8)])


# ── Save ──
prs.save(OUT_FILE)
print(f"Saved: {OUT_FILE}")
print(f"Slides: {len(prs.slides)}")
